import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0] # Title Layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Styling
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(15, 12, 8)
        
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(197, 160, 89)

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content Layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(197, 160, 89)
    
    text_frame = body_shape.text_frame
    text_frame.clear() # Clear default paragraph
    
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.name = 'Arial'
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(14)


def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, "Adaalat", "Digital Justice Platform\nEnd-to-End Workflow Presentation")
    
    # Slide 2
    add_content_slide(prs, "1. Vision & Overview", [
        "Problem: Finding the right lawyer and getting preliminary legal advice is often difficult, time-consuming, and intimidating.",
        "Solution: Adaalat bridges the gap between clients and legal professionals.",
        "Core Features:",
        "  • AI-Powered Legal Advisory (Initial assessment)",
        "  • Smart Lawyer Matching based on specialization",
        "  • Seamless Case Sharing (Smart Briefs)",
        "  • Secure Document Exchange"
    ])
    
    # Slide 3
    add_content_slide(prs, "2. User Onboarding & Authentication", [
        "Dual-Role Gateway:",
        "  • Users specify their purpose (Client vs. Lawyer) during Sign Up.",
        "Client Registration:",
        "  • Basic details (Name, Email, Phone).",
        "Lawyer Registration:",
        "  • Professional credentials (Bar Registration Number, Portfolio).",
        "Routing: Post-login, clients are routed to the Client Dashboard, while lawyers enter the dedicated Lawyer Dashboard."
    ])
    
    # Slide 4
    add_content_slide(prs, "3. Client Journey: AI Legal Advisory", [
        "Initial Consultation:",
        "  • Client submits a legal query via a chat-like interface.",
        "AI Analysis Engine:",
        "  • The backend (via RAG / GenAI) analyzes the query.",
        "Output Generation:",
        "  • Categorizes the legal area (e.g., Tenancy, Employment).",
        "  • Provides a preliminary analysis of the situation.",
        "  • Suggests actionable next steps.",
        "  • Automatically stores the query and analysis as a 'Case' in the backend."
    ])
    
    # Slide 5
    add_content_slide(prs, "4. Client Journey: Finding a Lawyer", [
        "Recommendation Engine:",
        "  • Based on the AI-determined legal area (e.g., Property Law), the system curates a list of specialized lawyers.",
        "Lawyer Profiles:",
        "  • Clients can view detailed portfolios, win rates, courts of practice, and consultation fees.",
        "Initiating Connection:",
        "  • Client clicks 'Connect'.",
        "  • A secure connection request containing the client's original query and AI analysis is sent to the specific lawyer."
    ])
    
    # Slide 6
    add_content_slide(prs, "5. Lawyer Dashboard: Smart Briefs", [
        "Lead Generation & Management:",
        "  • Lawyers receive the client's connection requests and AI advisories seamlessly in their 'Smart Briefs' feed.",
        "Pre-Analyzed Intelligence:",
        "  • Instead of raw queries, lawyers read structured briefs containing the categorized legal area, background, and recommended steps.",
        "Action Handling:",
        "  • Lawyers can choose to 'Accept' or 'Decline' the case directly from the dashboard."
    ])
    
    # Slide 7
    add_content_slide(prs, "6. Seamless Document Exchange", [
        "Post-Connection Workflow:",
        "  • Once accepted, the lawyer prepares necessary legal drafts, notices, or agreements.",
        "Lawyer Upload:",
        "  • The lawyer securely uploads PDFs targeted to specific clients.",
        "Client Receipt:",
        "  • The client dashboard updates instantly, notifying the user of incoming documents.",
        "  • Clients can view and download the PDFs directly inside Adaalat."
    ])
    
    # Slide 8
    add_content_slide(prs, "7. Technology Stack & Architecture", [
        "Frontend: React + Vite + Tailwind CSS",
        "  • Smooth animations, responsive glassmorphism UI.",
        "  • Separate entry points for distinct user experiences.",
        "Backend: Python Flask Server",
        "  • RESTful API for document management and case-state logic.",
        "  • Integration points ready for full LangChain/RAG pipeline.",
        "Data Storage:",
        "  • Local JSON-based caching for rapid prototyping (expandable to Postgres/Supabase)."
    ])

    # Slide 9
    add_content_slide(prs, "Summary", [
        "Adaalat transforms the fragmented legal process into a structured, user-friendly pipeline.",
        "1. Client asks AI ➞ 2. AI structures the problem ➞ 3. Client connects to Lawyer ➞ 4. Lawyer receives Smart Brief ➞ 5. Lawyer assists & shares documents.",
        "Thank You!"
    ])
    
    prs.save('Adaalat_Workflow.pptx')
    print("Presentation saved as Adaalat_Workflow.pptx")

if __name__ == '__main__':
    create_presentation()
